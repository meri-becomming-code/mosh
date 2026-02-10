
# Created by Meri Kasprak with the assistance of Gemini.
# Released freely under the GNU General Public License version 3. USE AT YOUR OWN RISK.

import os
import shutil
import re
from datetime import datetime
import zipfile
import base64
import uuid
import xml.etree.ElementTree as ET
import urllib.parse
from bs4 import BeautifulSoup


# --- Constants ---
ARCHIVE_FOLDER_NAME = "_ORIGINALS_DO_NOT_UPLOAD_"

def sanitize_filename(base_name):
    """
    Replaces spaces, dots, and special characters with underscores to ensure web safety.
    Input should be the filename WITHOUT extension.
    """
    # [STRICT FIX] Only allow letters, numbers, underscores, and hyphens. 
    # Everything else (including dots and commas) becomes an underscore.
    s_name = re.sub(r'[^\w\-]', '_', base_name)
    # Collapse multiple underscores
    s_name = re.sub(r'_+', '_', s_name)
    # Clean up trailing/leading underscores
    s_name = s_name.strip('_')
    return s_name

# --- Third Party Imports ---
try:
    import mammoth
except ImportError:
    mammoth = None

try:
    import openpyxl
    from openpyxl.utils import get_column_letter
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    Presentation = None

try:
    import fitz # PyMuPDF
except ImportError:
    fitz = None

try:
    import docx
except ImportError:
    docx = None

try:
    from pdfminer.high_level import extract_text
except ImportError:
    extract_text = None


# --- HTML Templates ---
HTML_TEMPLATE = """<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        body {{ 
            font-family: 'Segoe UI', 'Roboto', 'Helvetica Neue', Arial, sans-serif; 
            font-size: 16px;
            line-height: 1.6; 
            color: #333;
            max-width: none; 
            margin: 0; 
            padding: 0; 
            background-color: #f4f7f9;
        }}
        .main-content {{
            max-width: 1200px;
            margin: 0 auto;
            padding: 40px;
        }}
        h1 {{ 
            font-size: 2.25em; 
            font-weight: 700;
            color: #4b3190; 
            border-bottom: 2px solid #4b3190; 
            padding-bottom: 15px; 
            margin-bottom: 40px;
            text-align: center;
        }}
        h2 {{ color: #2c3e50; margin-top: 40px; font-weight: 600; border-bottom: 1px solid #dee2e6; padding-bottom: 5px; }}
        h3 {{ color: #444; margin-top: 30px; font-weight: 600; }}
        a {{ color: #0056b3; text-decoration: none; font-weight: 500; }}
        a:hover {{ text-decoration: underline; }}
        
        /* Table Styles */
        table {{ border-collapse: collapse; width: 100%; margin: 25px 0; font-size: 0.95em; border: 1px solid #ddd; background-color: #fff; }}
        th, td {{ border: 1px solid #ddd; padding: 12px 15px; text-align: left; }}
        th {{ background-color: #f8f9fa; font-weight: 600; color: #495057; }}
        tr:nth-child(even) {{ background-color: #f9f9f9; }}
        .content-table {{ width: 100%; border-collapse: collapse; }}
        
        img {{ max-width: 100%; height: auto; border-radius: 4px; border: 1px solid #eee; }}
        .grading-note {{ background-color: #e8f5e9; padding: 15px; border-left: 5px solid #4caf50; font-style: italic; border-radius: 0 4px 4px 0; }}
        .note {{ font-size: 0.95em; color: #555; background: #fff3cd; padding: 20px; border-radius: 8px; border: 1px solid #ffeeba; margin-bottom: 25px; }}
        
        /* Code Block Styles */
        code {{ background-color: #f1f3f5; padding: 2px 5px; border-radius: 4px; font-family: Consolas, 'Courier New', monospace; color: #d63384; }}
        pre {{ 
            background-color: #272822; 
            color: #f8f8f2; 
            padding: 15px; 
            border-radius: 8px; 
            overflow-x: auto; 
            font-family: Consolas, 'Courier New', monospace; 
            line-height: 1.4;
            margin: 20px 0;
            box-shadow: inset 0 2px 4px rgba(0,0,0,0.2);
        }}
        .code-block {{ margin: 15px 0; }}
        
        .slide-container {{ 
            overflow: auto; 
            clear: both; 
            margin-bottom: 60px; 
            padding: 60px; 
            border: 2px solid #ccc;
            border-top: 5px solid #4b3190;
            border-radius: 12px; 
            background-color: #fff;
            box-shadow: 0 8px 30px rgba(0,0,0,0.1);
            position: relative;
        }}
        .slide-container::after {{ content: ""; display: table; clear: both; }}
        .slide-title {{ margin-top: 0; padding-bottom: 10px; border-bottom: 1px solid #eee; margin-bottom: 25px; }}
        .slide-num {{ position: absolute; top: 15px; right: 25px; font-size: 0.8em; color: #666; font-weight: bold; }}
        .slide-image {{ border-radius: 6px; box-shadow: 0 4px 12px rgba(0,0,0,0.1); border: 1px solid #eee; }}
        
        /* Accounting & Excel Styles */
        .accounting-table {{ border-collapse: collapse; margin: 25px 0; font-family: 'Courier New', Courier, monospace; width: auto; min-width: 50%; }}
        .accounting-table th, .accounting-table td {{ border: 1px solid #ccc; padding: 8px 12px; }}
        .currency-cell {{ text-align: right; white-space: nowrap; }}
        .label-cell {{ text-align: left; }}
        .total-row {{ font-weight: bold; border-top: 2px solid #333; }}
        .grand-total {{ border-bottom: 3px double #333; }}
        .negative {{ color: #d32f2f; }}
        .excel-sheet-header {{ 
            background-color: #1f6e43; 
            color: white; 
            padding: 10px 20px; 
            margin-top: 40px; 
            border-radius: 4px 4px 0 0;
            display: inline-block;
        }}
        .excel-container {{ overflow-x: auto; margin-bottom: 50px; }}
        
        /* Dynamic Style Overrides */
        {style_overrides}
    </style>
</head>
<body>
    <div class="main-content" style="max-width: 1200px; margin: 0 auto; padding: 40px;">
        <h1>{title}</h1>
        {content}
    </div>
</body>
</html>

"""


def _save_html(content, title, source_file, output_path, style_overrides=""):
    """Wraps content in template and saves file."""
    html = HTML_TEMPLATE.format(
        title=title,
        content=content,
        style_overrides=style_overrides
    )
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(html)
    return output_path

def optimize_image(image_path, max_width=1100, make_transparent=False):
    """Resizes, compresses, and optionally removes white backgrounds from images."""
    try:
        from PIL import Image, ImageDraw
        img = Image.open(image_path).convert("RGBA")
        
        # 1. Resize if too wide (prevent Canvas bloat)
        w, h = img.size
        if w > max_width:
            ratio = max_width / float(w)
            new_h = int(float(h) * ratio)
            img = img.resize((max_width, new_h), Image.Resampling.LANCZOS)
            w, h = max_width, new_h

        # 2. Magic Background Removal (Optional)
        if make_transparent:
            # Check corners for white-ish color
            for corner in [(0,0), (w-1, 0), (0, h-1), (w-1, h-1)]:
                # If corner is purely white, floodfill transparency
                p = img.getpixel(corner)
                if p[0] > 240 and p[1] > 240 and p[2] > 240:
                    ImageDraw.floodfill(img, xy=corner, value=(255, 255, 255, 0), thresh=15)
        
        # 3. Save optimized
        img.save(image_path, "PNG", optimize=True)
        return True
    except Exception as e:
        print(f"Image Optimization failed for {image_path}: {e}")
        return False

def extract_theme_info(prs):
    """
    Extracts theme colors and fonts from a PowerPoint presentation.
    Returns a dict with color scheme and potentially font names.
    """
    theme_info = {'colors': {}, 'font': 'inherit'}
    try:
        # Access the theme part
        # [NOTE] Presentation.part.package.part_related_by is a common way to get specific related parts
        theme_rel_type = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        theme_part = prs.part.package.part_related_by(theme_rel_type)
        if not theme_part: return theme_info
        
        xml_content = theme_part.blob
        root = ET.fromstring(xml_content)
        
        # Namespaces
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        
        # 1. Extract Colors
        clr_scheme = root.find('.//a:clrScheme', ns)
        if clr_scheme is not None:
            # Common mapping for Office colors
            mapping = {
                'dk1': 'dark1', 'lt1': 'light1', 
                'accent1': 'accent1', 'accent2': 'accent2', 
                'accent3': 'accent3', 'accent4': 'accent4'
            }
            for tag, label in mapping.items():
                elem = clr_scheme.find(f'a:{tag}', ns)
                if elem is not None:
                    # Look for srgbClr (RGB)
                    srgb = elem.find('.//a:srgbClr', ns)
                    if srgb is not None:
                        val = srgb.get('val')
                        if val: theme_info['colors'][label] = f"#{val}"
                    # Or sysClr (System - often white/black)
                    else:
                        sys = elem.find('.//a:sysClr', ns)
                        if sys is not None:
                            last_clr = sys.get('lastClr')
                            if last_clr: theme_info['colors'][label] = f"#{last_clr}"
        
        # 2. Extract Fonts
        font_scheme = root.find('.//a:fontScheme', ns)
        if font_scheme is not None:
            # Prefer minorFont (body text)
            minor_font = font_scheme.find('.//a:minorFont/a:latin', ns)
            if minor_font is not None:
                typeface = minor_font.get('typeface')
                if typeface: theme_info['font'] = typeface
                    
    except Exception as e:
        # Silently fail for themes, it's a "nice to have"
        pass
    return theme_info

# --- Converters ---

def convert_docx_to_html(docx_path, io_handler=None):
    """Converts DOCX to HTML using Mammoth (with style mapping)."""
    if not mammoth:
        return None, "Mammoth library not installed."

    try:
        # Custom Style Map for Canvas alignment
        style_map = """
        p[style-name='Title'] => h1:fresh
        p[style-name='Heading 1'] => h2:fresh
        p[style-name='Heading 2'] => h3:fresh
        p[style-name='Heading 3'] => h4:fresh
        p[style-name='Heading 4'] => h5:fresh
        p[style-name='Complete/Incomplete'] => p.grading-note:fresh
        r[style-name='Strong'] => strong
        b => strong
        i => em
        """
        
        filename = os.path.splitext(os.path.basename(docx_path))[0]
        output_dir = os.path.dirname(docx_path)
        
        # Image Handler for Mammoth
        def convert_image(image):
            # 1. Create web_resources/[filename] folder if it doesn't exist
            safe_filename = sanitize_filename(filename)
            res_dir = os.path.join(output_dir, "web_resources", safe_filename)
            if not os.path.exists(res_dir):
                os.makedirs(res_dir)
            
            # 2. Extract description (from original doc)
            original_alt = image.alt_text if image.alt_text else f"Image from {filename}"
            
            # 3. Save Image File
            with image.open() as image_source:
                image_bytes = image_source.read()
            
            # Generate unique name
            ext_img = image.content_type.split('/')[-1]
            if ext_img == 'jpeg': ext_img = 'jpg'
            short_id = uuid.uuid4().hex[:8]
            img_name = f"img_{short_id}.{ext_img}"
            img_path = os.path.join(res_dir, img_name)
            
            with open(img_path, 'wb') as f:
                f.write(image_bytes)
            
            # [ENHANCED] Get Natural Dimensions via Pillow
            from PIL import Image as PILImage
            import io
            width_attr = "100%"
            style_attr = "max-width: 600px; height: auto;" # Safe default
            
            try:
                with PILImage.open(io.BytesIO(image_bytes)) as pil_img:
                    w, h = pil_img.size
                    if w < 200:
                        width_attr = str(w)
                        style_attr = "" # Keep natural
                    else:
                        width_attr = str(min(w, 800))
            except: pass

            # [INTERACTIVE] Prompt for Alt Text
            final_alt = original_alt
            if io_handler:
                import interactive_fixer
                mem_key = interactive_fixer.normalize_image_key(img_name, img_path)
                if mem_key in io_handler.memory:
                    final_alt = io_handler.memory[mem_key]
                else:
                    choice = io_handler.prompt_image(f"   > Alt Text for {img_name} (or Enter to keep original): ", img_path, context=f"Context: {filename} (Word Document)").strip()
                    if choice:
                        if choice == "__DECORATIVE__":
                            final_alt = ""
                        else:
                            final_alt = choice
                        io_handler.memory[mem_key] = final_alt
                        io_handler.save_memory()

            # 4. Return Tag with Standard Relative Path
            return {
                "src": f"web_resources/{safe_filename}/{img_name}",
                "alt": final_alt,
                "width": width_attr,
                "style": style_attr
            }

        with open(docx_path, "rb") as docx_file:
            result = mammoth.convert_to_html(docx_file, style_map=style_map, convert_image=mammoth.images.img_element(convert_image))
            html_content = result.value
            messages = result.messages # Warnings
            
        # Logging for user visibility
        img_count = html_content.count("<img")
        print(f"    [LOG] Extracted {img_count} images from Word document.")

        # Post-Processing: Basic Cleanup
        # Remove empty paragraphs often generated by extra Returns in Word
        html_content = html_content.replace("<p></p>", "").replace("<p>&nbsp;</p>", "")
        
        # Ensure tables have some basic class for our CSS
        html_content = html_content.replace("<table>", '<table class="content-table">')
        
        # [NEW] Remove empty tables that often come from Word formatting
        temp_soup = BeautifulSoup(html_content, 'html.parser')
        tables_removed = 0
        for table in temp_soup.find_all('table'):
            has_content = False
            for cell in table.find_all(['td', 'th']):
                if cell.get_text(strip=True) or cell.find('img'):
                    has_content = True
                    break
            if not has_content:
                table.extract()
                tables_removed += 1
        
        if tables_removed > 0:
            print(f"    [LOG] Removed {tables_removed} empty tables from Word document.")
            html_content = str(temp_soup)

        s_filename = sanitize_filename(filename)
        output_path = os.path.join(output_dir, f"{s_filename}.html")
        
        # Wrap in template
        _save_html(html_content, filename, docx_path, output_path)
        
        return output_path, None
    except Exception as e:
        return None, str(e)


def convert_excel_to_html(xlsx_path):
    """Converts Excel to HTML Tables using OpenPyXL."""
    if not openpyxl:
        return None, "OpenPyXL library not installed."

    try:
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
        html_parts = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            html_parts.append(f'<div class="excel-container">')
            html_parts.append(f'<h3 class="excel-sheet-header">Sheet: {sheet_name}</h3>')
            html_parts.append('<table class="accounting-table">')
            
            # [ACCOUNTING FIX] Detect merged cells to handle headers correctly if possible
            # (Basic implementation: just treat them as individual cells for now to avoid complexity)
            
            rows = list(ws.rows)
            if rows:
                # 1. Header Row
                html_parts.append("<thead><tr>")
                for cell in rows[0]:
                    val = cell.value if cell.value is not None else ""
                    # Use th with scope for accessibility
                    html_parts.append(f'<th scope="col">{val}</th>')
                html_parts.append("</tr></thead>")
                
                # 2. Body Rows
                html_parts.append("<tbody>")
                for row in rows[1:]:
                    html_parts.append("<tr>")
                    for cell in row:
                        val = cell.value
                        if val is None:
                            html_parts.append("<td></td>")
                            continue
                            
                        # Detection: Style Classes
                        classes = []
                        
                        # A. Alignment (Numbers right, Text left)
                        if isinstance(val, (int, float, datetime)):
                            classes.append("currency-cell")
                        else:
                            classes.append("label-cell")
                            
                        # B. Number Formatting (Currency, Percent, Accounting)
                        fmt = cell.number_format
                        str_val = str(val)
                        
                        if fmt:
                            if "$" in fmt or "Currency" in fmt or "Accounting" in fmt:
                                try:
                                    str_val = f"${val:,.2f}"
                                    if val < 0: 
                                        classes.append("negative")
                                        # Accounting format often uses ( ) for negatives
                                        str_val = f"({str_val.replace('-', '')})"
                                except: pass
                            elif "%" in fmt:
                                try: str_val = f"{val*100:.1f}%"
                                except: pass
                            elif "yyyy" in fmt or "mm" in fmt:
                                try: str_val = val.strftime("%Y-%m-%d")
                                except: pass
                        
                        # C. Borders (Total Rows)
                        if cell.border:
                            if cell.border.bottom and cell.border.bottom.style:
                                if cell.border.bottom.style == 'double':
                                    classes.append("grand-total")
                                else:
                                    classes.append("total-row")
                        
                        # D. Font (Bold)
                        if cell.font and cell.font.bold:
                            classes.append("total-row") # Use same bolding style

                        class_attr = f' class="{" ".join(classes)}"' if classes else ""
                        html_parts.append(f'<td{class_attr}>{str_val}</td>')
                        
                    html_parts.append("</tr>")
                html_parts.append("</tbody>")
            
            html_parts.append("</table>")
            html_parts.append("</div>") # End excel-container

        full_content = "\n".join(html_parts)
        
        filename = os.path.splitext(os.path.basename(xlsx_path))[0]
        s_filename = sanitize_filename(filename)
        output_path = os.path.join(os.path.dirname(xlsx_path), f"{s_filename}.html")
        
        _save_html(full_content, filename, xlsx_path, output_path)
        return output_path, None
        
    except Exception as e:
        return None, str(e)


def get_shape_text_styles(shape, theme=None):
    """Extracts CSS styles (color, background-color, border, rotation) from a shape."""
    styles = []
    
    # 1. Background Color (Shape Fill)
    try:
        if shape.fill.type == 1: # Solid fill
            rgb = shape.fill.fore_color.rgb
            if rgb: styles.append(f"background-color: #{rgb};")
    except: pass
    
    # 2. Border / Line
    try:
        if shape.line.fill.type == 1: # Solid line
            rgb = shape.line.color.rgb
            width = int(shape.line.width / 12700) # CM to Px approx
            if rgb: styles.append(f"border: {width}px solid #{rgb};")
    except: pass

    # 3. Rotation
    try:
        if shape.rotation != 0:
            styles.append(f"transform: rotate({shape.rotation}deg);")
    except: pass

    # 4. Text Color (Looking at first paragraph/run)
    try:
        if shape.has_text_frame and shape.text_frame.paragraphs:
            para = shape.text_frame.paragraphs[0]
            if para.runs:
                rgb = para.runs[0].font.color.rgb
                if rgb: styles.append(f"color: #{rgb};")
            elif para.font.color and para.font.color.rgb:
                 styles.append(f"color: #{para.font.color.rgb};")
    except: pass
    
    # 5. Padding/Border if background/border exists
    if any(k in s for s in styles for k in ["background-color", "border"]):
        styles.append("padding: 15px; border-radius: 8px; margin-bottom: 10px;")
    
    return " ".join(styles)

def get_image_styles(shape):
    """Specific styles for pictures like borders and rotation."""
    styles = []
    try:
        if shape.rotation != 0:
            styles.append(f"transform: rotate({shape.rotation}deg);")
        
        if shape.line.fill.type == 1:
            rgb = shape.line.color.rgb
            width = int(shape.line.width / 12700)
            if rgb: styles.append(f"border: {width}px solid #{rgb};")
    except: pass
    return " ".join(styles)


def extract_all_shapes_recursive(shapes):
    """Recursively flattened list of shapes (handles groups)."""
    flat_list = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            flat_list.extend(extract_all_shapes_recursive(shape.shapes))
        else:
            flat_list.append(shape)
    return flat_list


def convert_ppt_to_html(ppt_path, io_handler=None):
    """Converts PPTX to HTML Lecture Notes + Extracts Images."""
    if not Presentation:
        return None, "python-pptx library not installed."

    try:
        prs = Presentation(ppt_path)
        filename = os.path.splitext(os.path.basename(ppt_path))[0]
        output_dir = os.path.dirname(ppt_path)
        
        # [THEME AWARENESS] Extract theme data
        theme = extract_theme_info(prs)
        style_overrides = ""
        
        # Apply Body Font if found
        if theme['font'] != 'inherit':
            style_overrides += f"body {{ font-family: '{theme['font']}', sans-serif; }}\n"
            
        # Apply Accent Colors to Slide Containers and Headings
        accent1 = theme['colors'].get('accent1', '#4b3190') # Default purple if not found
        dark1 = theme['colors'].get('dark1', '#333')
        light1 = theme['colors'].get('light1', '#fff')
        
        style_overrides += f"""
            .slide-container {{ border-top-color: {accent1}; border-top-width: 5px; border-left: 2px solid #ccc; border-right: 2px solid #ccc; border-bottom: 2px solid #ccc; background-color: {light1}; }}
            .slide-title {{ color: {dark1}; border-bottom-color: {accent1}; }}
            h1 {{ color: {accent1}; border-bottom-color: {accent1}; }}
            h2 {{ border-bottom-color: {accent1}; }}
        """


        html_parts = []
        
        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            
            # [NEW] Inline style for slide container (Canvas survival)
            slide_style = (
                f"margin-bottom: 60px; padding: 60px; border: 2px solid #ccc; "
                f"border-top: 5px solid {accent1}; border-radius: 12px; "
                f"background-color: {light1}; box-shadow: 0 8px 30px rgba(0,0,0,0.1); "
                f"position: relative; overflow: auto; clear: both;"
            )
            html_parts.append(f'<div class="slide-container" id="slide-{slide_num}" style="{slide_style}">')
            html_parts.append(f'<div class="slide-num" style="position: absolute; top: 15px; right: 25px; font-size: 0.8em; color: #666; font-weight: bold;">Slide {slide_num}</div>')

            
            # [NEW] Detect if slide has text content (for image sizing)
            has_text_content = False
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    if shape.text_frame.text.strip():
                        has_text_content = True
                        break
            
            # Title
            if slide.shapes.title:
                title_text = slide.shapes.title.text_frame.text
                html_parts.append(f'<h2 class="slide-title">{title_text}</h2>')
            
            # Content (Text & Images)
            # [BARNEY FIX] Recursive extraction to catch text inside Groups
            all_shapes = extract_all_shapes_recursive(slide.shapes)
            
            # Sort shapes by position to ensure reading order and side-by-side floating
            # We round to nearest 10 pixels to group items that are roughly on the same vertical line
            def shape_sort_key(s):
                try:
                    top = getattr(s, 'top', 0)
                    left = getattr(s, 'left', 0)
                    # Prioritize images slightly if they are on the same line to ensure they float correctly
                    priority = 0 if s.shape_type == MSO_SHAPE_TYPE.PICTURE else 1
                    return (round(top / 95250) * 10, priority, left)
                except:
                    return (0, 0, 0)

            sorted_shapes = sorted(all_shapes, key=shape_sort_key)
            
            for shape in sorted_shapes:
                # Text
                if shape.has_text_frame:
                    if shape == slide.shapes.title: continue 

                    # [SMART FIX] 1. Code Block Detection (Monospace Fonts)
                    is_code = False
                    bg_color = None
                    text_color = None
                    
                    if shape.text_frame.paragraphs:
                        # Check first paragraph font
                        para = shape.text_frame.paragraphs[0]
                        font_name = para.font.name
                        if font_name and any(f in font_name.lower() for f in ['courier', 'consolas', 'mono', 'lucida console']):
                            is_code = True
                            
                            # Try to extract Background Color from Shape Fill
                            try:
                                if shape.fill.type == 1: # Solid fill
                                    rgb = shape.fill.fore_color.rgb
                                    bg_color = f"#{rgb}"
                            except: pass
                            
                            # Try to extract Text Color from first run
                            try:
                                if para.runs:
                                    rgb = para.runs[0].font.color.rgb
                                    if rgb:
                                        text_color = f"#{rgb}"
                            except: pass

                    if is_code:
                        safe_text = shape.text_frame.text.replace("<", "&lt;").replace(">", "&gt;")
                        style = ""
                        if bg_color: style += f"background-color: {bg_color}; "
                        if text_color: style += f"color: {text_color}; "
                        
                        html_parts.append(f'<pre class="code-block" style="{style}">{safe_text}</pre>')
                        continue

                    # [NEW] Extract Text Box Styles (Colors/Backgrounds)
                    box_style = get_shape_text_styles(shape, theme)
                    if box_style:
                        html_parts.append(f'<div class="text-box" style="{box_style}">')

                    # [SMART FIX] 2. Improved Bullet Detection + Hyperlink Preservation
                    text_content = []
                    for paragraph in shape.text_frame.paragraphs:
                        if not paragraph.text.strip(): continue
                        
                        # Build paragraph content from runs to preserve links
                        para_html_parts = []
                        for run in paragraph.runs:
                            run_text = run.text.replace("<", "&lt;").replace(">", "&gt;")
                            if not run_text: continue
                            
                            # Check for Hyperlink
                            hlink = run.hyperlink.address
                            if hlink:
                                para_html_parts.append(f'<a href="{hlink}">{run_text}</a>')
                            else:
                                # Preserving Styles (Bold, Italic, Color, Font)
                                transformed = run_text
                                try:
                                    inline_styles = []
                                    if run.font.color and run.font.color.rgb:
                                        inline_styles.append(f"color: #{run.font.color.rgb};")
                                    if run.font.name:
                                        inline_styles.append(f"font-family: '{run.font.name}', sans-serif;")
                                    if run.font.size:
                                        size_pt = int(run.font.size / 12700)
                                        # Only keep if >= 10pt per user request for readability
                                        if size_pt >= 10:
                                            inline_styles.append(f"font-size: {size_pt}pt;")
                                    
                                    if inline_styles:
                                        transformed = f'<span style="{" ".join(inline_styles)}">{transformed}</span>'
                                    
                                    if run.font.bold: transformed = f"<strong>{transformed}</strong>"
                                    if run.font.italic: transformed = f"<em>{transformed}</em>"
                                except: pass
                                para_html_parts.append(transformed)
                        
                        full_para_html = "".join(para_html_parts)
                        if not full_para_html.strip(): continue

                        # Check if this paragraph is actually a bullet
                        is_bullet = False
                        try:
                            if paragraph.level > 0:
                                is_bullet = True
                            elif paragraph.text.strip().startswith(('•', '-', '*', '◦', '▪')):
                                is_bullet = True
                        except: pass

                        if is_bullet:
                            text_content.append(f"<li>{full_para_html}</li>")
                        else:
                            text_content.append(f"<p>{full_para_html}</p>")
                    
                    if text_content:
                        final_shape_html = ""
                        in_list = False
                        for item in text_content:
                            if item.startswith("<li>"):
                                if not in_list:
                                    final_shape_html += "<ul>"
                                    in_list = True
                                final_shape_html += item
                            else:
                                if in_list:
                                    final_shape_html += "</ul>"
                                    in_list = False
                                final_shape_html += item
                        if in_list:
                            final_shape_html += "</ul>"
                        html_parts.append(final_shape_html)
                    
                    if box_style:
                        html_parts.append('</div>')

                # Tables
                if shape.has_table:
                    # [NEW] Check if table is empty
                    is_empty = True
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame and cell.text_frame.text.strip():
                                is_empty = False
                                break
                    
                    if not is_empty:
                        html_parts.append('<table class="content-table" border="1">')
                        for row in shape.table.rows:
                            html_parts.append('<tr>')
                            for cell in row.cells:
                                # Extract text from cell
                                cell_text = ""
                                if cell.text_frame:
                                    cell_text = cell.text_frame.text.strip()
                                html_parts.append(f'<td>{cell_text}</td>')
                            html_parts.append('</tr>')
                        html_parts.append('</table>')

                # Images (Alt Text prompts only if no Silent Memory)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        ext = image.ext
                        safe_filename = sanitize_filename(filename)
                        res_dir = os.path.join(output_dir, "web_resources", safe_filename)
                        if not os.path.exists(res_dir): os.makedirs(res_dir)
                        
                        image_filename = f"slide{slide_num}_{uuid.uuid4().hex[:6]}.{ext}"
                        image_full_path = os.path.join(res_dir, image_filename)
                        
                        # 1. Save original bytes first
                        with open(image_full_path, 'wb') as img_f:
                            img_f.write(image_bytes)
                        
                        # 2. [NEW] Image Optimization & Magic Transparency
                        # We save as PNG for transparency support
                        optimize_image(image_full_path, max_width=400, make_transparent=True)
                        
                        rel_path = f"web_resources/{safe_filename}/{image_filename}"
                        width_px = int(shape.width / 9525) if hasattr(shape, 'width') else 400
                        
                        if has_text_content:
                            max_width = 450
                            if width_px > max_width: width_px = max_width
                        else:
                            if width_px > 800: width_px = 800
                        
                        slide_width = prs.slide_width if hasattr(prs, 'slide_width') else 9144000
                        shape_left = shape.left if hasattr(shape, 'left') else 0
                        shape_center_x = shape_left + (shape.width / 2) if hasattr(shape, 'width') else 0
                        
                        center_threshold = slide_width * 0.1
                        dist_from_center = abs(shape_center_x - (slide_width / 2))
                        
                        if dist_from_center < center_threshold:
                            float_style = "display: block; margin: 20px auto;"
                        elif shape_center_x < slide_width / 2:
                            float_style = "float: left; margin: 0 20px 15px 0;"
                        else:
                            float_style = "float: right; margin: 0 0 15px 20px;"
                        
                        # [NEW] Enhanced Image Styles (Borders/Rotation)
                        extra_img_style = get_image_styles(shape)
                        final_img_style = f"{float_style} {extra_img_style}".strip()
                        
                        # [SMART FIX] Silent Memory and prompt
                        alt_text = "" # Default to decorative/empty if skipped
                        if io_handler:
                            import interactive_fixer
                            mem_key = interactive_fixer.normalize_image_key(rel_path, image_full_path)
                            
                            if mem_key in io_handler.memory:
                                alt_text = io_handler.memory[mem_key]
                            else:
                                slide_title = slide.shapes.title.text_frame.text if slide.shapes.title else f"Slide {slide_num}"
                                choice = io_handler.prompt_image(f"   > Alt Text for Slide {slide_num} image (or Enter to skip): ", image_full_path, context=f"Context: {slide_title}").strip()
                                if choice:
                                    if choice == "__DECORATIVE__":
                                        alt_text = ""
                                    else:
                                        alt_text = choice
                                    io_handler.memory[mem_key] = alt_text
                                    io_handler.save_memory()

                        html_parts.append(f'<img src="{rel_path}" alt="{alt_text}" width="{width_px}" class="slide-image" style="{final_img_style}">')
                    except Exception as img_err:
                        print(f"Skipped image on slide {slide_num}: {img_err}")

            # [NEW] Capture Speaker Notes (Essential context often missed)
            try:
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        html_parts.append('<div class="speaker-notes" style="margin-top: 30px; padding: 20px; background: #f9f9f9; border-left: 4px solid #4b3190; font-style: italic;">')
                        html_parts.append(f'<strong>Speaker Notes:</strong><br>{notes_text.replace("\n", "<br>")}')
                        html_parts.append('</div>')
            except: pass

            html_parts.append('</div>')

        full_content = "\n".join(html_parts)
        s_filename = sanitize_filename(filename)
        output_path = os.path.join(output_dir, f"{s_filename}.html")
        
        _save_html(full_content, filename, ppt_path, output_path, style_overrides=style_overrides)
        return output_path, None

    except Exception as e:
        return None, str(e)


def convert_pdf_to_html(pdf_path, io_handler=None):
    """Converts PDF to HTML using PyMuPDF (Images + Text)."""
    if not fitz:
        if not extract_text:
             return None, "Neither PyMuPDF (fitz) nor pdfminer.six are installed."
        else:
             # Fallback to old method if fitz is missing (though we just installed it)
             return _convert_pdf_fallback(pdf_path)

    try:
        doc = fitz.open(pdf_path)
        filename = os.path.splitext(os.path.basename(pdf_path))[0]
        output_dir = os.path.dirname(pdf_path)
        
        html_parts = []
        html_parts.append('<div class="pdf-content">')
        
        total_text_blocks = 0
        for i, page in enumerate(doc):
            page_num = i + 1
            html_parts.append(f'<div class="page-container" id="page-{page_num}" style="margin-bottom: 30px; border-bottom: 1px solid #ccc; padding-bottom: 20px;">')
            html_parts.append(f'<p class="note" style="font-size: 0.8em; color: #666; font-weight: bold;">Page {page_num}</p>')
            
            # [IMPROVED] Extract tables FIRST to know their positions
            table_regions = []
            try:
                tables = page.find_tables()
                if tables.tables:
                    for tab in tables:
                        # Get the bounding box of the table
                        bbox = tab.bbox  # (x0, y0, x1, y1)
                        table_regions.append({
                            'bbox': bbox,
                            'table': tab
                        })
            except Exception as e:
                print(f"Table detection failed on page {page_num}: {e}")
            
            # 1. Extract Content via Dict (Structure + Images)
            page_dict = page.get_text("dict")
            blocks = page_dict.get("blocks", [])
            
            # Helper function to check if a position overlaps with any table
            def get_table_at_position(y_pos):
                for tr in table_regions:
                    bbox = tr['bbox']
                    # Check if y_pos is within table's vertical range
                    if bbox[1] <= y_pos <= bbox[3]:
                        return tr
                return None
            
            # Track which tables we've already inserted
            inserted_tables = set()
            
            for block in blocks:
                # Check if we should insert a table before this block
                block_top = block['bbox'][1]
                table_region = get_table_at_position(block_top)
                
                if table_region and id(table_region) not in inserted_tables:
                    # Insert table before this block
                    try:
                        tab = table_region['table']
                        # Extract table data
                        df = tab.to_pandas()
                        
                        # Build HTML table with proper structure
                        html_parts.append('\u003ctable class="content-table"\u003e')
                        
                        # Detect if first row looks like headers (all strings, capitalized, etc.)
                        first_row = df.iloc[0] if len(df) > 0 else []
                        is_header = all(isinstance(v, str) for v in first_row if v is not None)
                        
                        if is_header and len(df) > 1:
                            # Use first row as header
                            html_parts.append('\u003cthead\u003e\u003ctr\u003e')
                            for cell in first_row:
                                html_parts.append(f'\u003cth\u003e{cell if cell else ""}\u003c/th\u003e')
                            html_parts.append('\u003c/tr\u003e\u003c/thead\u003e')
                            # Rest as body
                            html_parts.append('\u003ctbody\u003e')
                            for _, row in df.iloc[1:].iterrows():
                                html_parts.append('\u003ctr\u003e')
                                for cell in row:
                                    html_parts.append(f'\u003ctd\u003e{cell if cell else ""}\u003c/td\u003e')
                                html_parts.append('\u003c/tr\u003e')
                            html_parts.append('\u003c/tbody\u003e')
                        else:
                            # No headers, all rows as data
                            html_parts.append('\u003ctbody\u003e')
                            for _, row in df.iterrows():
                                html_parts.append('\u003ctr\u003e')
                                for cell in row:
                                    html_parts.append(f'\u003ctd\u003e{cell if cell else ""}\u003c/td\u003e')
                                html_parts.append('\u003c/tr\u003e')
                            html_parts.append('\u003c/tbody\u003e')
                        
                        html_parts.append('\u003c/table\u003e')
                        inserted_tables.add(id(table_region))
                    except Exception as e:
                        print(f"Error rendering table: {e}")
                
                # Skip blocks that are part of tables
                block_bbox = block['bbox']
                is_in_table = False
                for tr in table_regions:
                    tab_bbox = tr['bbox']
                    # Check if block is completely within table bounds
                    if (tab_bbox[0] <= block_bbox[0] and block_bbox[2] <= tab_bbox[2] and 
                        tab_bbox[1] <= block_bbox[1] and block_bbox[3] <= tab_bbox[3]):
                        is_in_table = True
                        break
                
                if is_in_table:
                    continue
                
                # Type 1 = Image
                if block['type'] == 1:
                    try:
                        ext = block['ext']
                        image_bytes = block['image']
                        
                        # [SIZE FIX] Get dimensions from BBox (in points)
                        bbox = block['bbox'] # (x0, y0, x1, y1)
                        width_pt = bbox[2] - bbox[0]
                        width_attr = int(width_pt)
                        
                        # [NEW] Alignment Detection
                        # PDF page width (usually ~600pt)
                        page_width = page.rect.width
                        shape_center_x = bbox[0] + (width_pt / 2)
                        
                        # Alignment logic
                        center_threshold = page_width * 0.1
                        dist_from_center = abs(shape_center_x - (page_width / 2))
                        
                        if dist_from_center < center_threshold:
                            float_style = "display: block; margin: 20px auto;"
                        elif shape_center_x < page_width / 2:
                            float_style = "float: left; margin: 0 20px 15px 0;"
                        else:
                            float_style = "float: right; margin: 0 0 15px 20px;"

                        # Save Image
                        safe_filename = sanitize_filename(filename)
                        res_dir = os.path.join(output_dir, "web_resources", safe_filename)
                        if not os.path.exists(res_dir): os.makedirs(res_dir)

                        image_filename = f"page{page_num}_img_{uuid.uuid4().hex[:6]}.{ext}"
                        image_full_path = os.path.join(res_dir, image_filename)
                        
                        with open(image_full_path, "wb") as f:
                            f.write(image_bytes)
                            
                        rel_path = f"web_resources/{safe_filename}/{image_filename}"
                        
                        # [INTERACTIVE] Prompt for Alt Text
                        alt_text = f"Image from Page {page_num}"
                        if io_handler:
                            import interactive_fixer
                            mem_key = interactive_fixer.normalize_image_key(rel_path, image_full_path)
                            if mem_key in io_handler.memory:
                                alt_text = io_handler.memory[mem_key]
                            else:
                                choice = io_handler.prompt_image(f"   > Alt Text for Page {page_num} image (or Enter to skip): ", image_full_path, context=f"Context: PDF Page {page_num}").strip()
                                if choice:
                                    if choice == "__DECORATIVE__":
                                        alt_text = ""
                                    else:
                                        alt_text = choice
                                    io_handler.memory[mem_key] = alt_text
                                    io_handler.save_memory()

                        html_parts.append(f'\u003cimg src="{rel_path}" alt="{alt_text}" width="{width_attr}" class="content-image" style="{float_style}"\u003e')
                    except Exception as e:
                        print(f"Skipped PDF image: {e}")

                # Type 0 = Text
                elif block['type'] == 0:
                     total_text_blocks += 1
                     # [IMPROVED] Aggregate all text in this block first, then intelligently group
                     block_lines = []
                     
                     for line in block["lines"]:
                         # Combine all spans in this line into a single text + metadata
                         line_text = ""
                         line_font_size = 0
                         line_y_pos = line["bbox"][1]  # Top Y coordinate
                         
                         for span in line["spans"]:
                             text = span["text"].strip()
                             if text:
                                 if line_text:
                                     line_text += " "
                                 line_text += text
                                 line_font_size = max(line_font_size, span["size"])
                         
                         if line_text:
                             block_lines.append({
                                 'text': line_text,
                                 'font_size': line_font_size,
                                 'y_pos': line_y_pos
                             })
                     
                     if not block_lines:
                         continue
                     
                     # Now group lines into semantic units (paragraphs, lists, headers)
                     i = 0
                     while i < len(block_lines):
                         current_line = block_lines[i]
                         text = current_line['text']
                         font_size = current_line['font_size']
                         safe_text = text.replace("<", "&lt;").replace(">", "&gt;")
                         
                         # Check for bullets first (priority over headers)
                         is_bullet = text.startswith(('• ', '- ', '* ', '◦ ', '▪ ', '⚬ '))
                         
                         if is_bullet:
                             # Collect consecutive bullet points
                             html_parts.append("<ul>")
                             while i < len(block_lines) and block_lines[i]['text'].startswith(('• ', '- ', '* ', '◦ ', '▪ ', '⚬ ')):
                                 item_text = block_lines[i]['text'].replace("<", "&lt;").replace(">", "&gt;")
                                 html_parts.append(f"<li>{item_text}</li>")
                                 i += 1
                             html_parts.append("</ul>")
                             continue
                         
                         # Check if header
                         if font_size > 18:
                             html_parts.append(f"<h2>{safe_text}</h2>")
                             i += 1
                             continue
                         elif font_size > 14:
                             html_parts.append(f"<h3>{safe_text}</h3>")
                             i += 1
                             continue
                         
                         # Otherwise, group into paragraph
                         paragraph_lines = [safe_text]
                         i += 1
                         
                         while i < len(block_lines):
                             next_line = block_lines[i]
                             if next_line['font_size'] > 14:  # Next is a header
                                 break
                             if next_line['text'].startswith(('• ', '- ', '* ')):  # Next is a bullet
                                 break
                             
                             # Check vertical gap
                             y_gap = abs(next_line['y_pos'] - block_lines[i-1]['y_pos'])
                             if y_gap > 24:  # Large gap = new paragraph
                                 break
                             
                             paragraph_lines.append(next_line['text'].replace("<", "&lt;").replace(">", "&gt;"))
                             i += 1
                         
                         html_parts.append(f"<p>{' '.join(paragraph_lines)}</p>")
            
            # [NEW] Fallback Image Extraction (Catch missed XObjects)
            try:
                img_list = page.get_images(full=True)
                # Filter out images already found in the dict block loop
                # (Simple heuristic: check if we've already saved images for this page in rel_path)
                found_count = len([p for p in html_parts if "web_resources" in p and f"page{page_num}_img" in p])
                
                if len(img_list) > found_count:
                    self.gui_handler.log(f"   [PDF] Pass 2: Found {len(img_list) - found_count} additional images...")
                    for img_index, img in enumerate(img_list):
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        ext = base_image["ext"]
                        
                        # Use a unique name for fallback images
                        fallback_name = f"page{page_num}_fallback_{xref}.{ext}"
                        fallback_path = os.path.join(res_dir, fallback_name)
                        
                        if not os.path.exists(fallback_path):
                            with open(fallback_path, "wb") as f:
                                f.write(image_bytes)
                                
                            rel_path = f"web_resources/{safe_filename}/{fallback_name}"
                            alt_text = f"Fallback Image from Page {page_num}"
                            
                            # Prompt for alt text if possible
                            if io_handler:
                                choice = io_handler.prompt_image(f"   > Missing Alt Text for PDF image on Page {page_num} (from fallback): ", fallback_path, context=f"Context: PDF Page {page_num}").strip()
                                if choice:
                                    alt_text = choice if choice != "__DECORATIVE__" else ""

                            html_parts.append(f'\u003cimg src="{rel_path}" alt="{alt_text}" class="content-image" style="display: block; margin: 20px auto; max-width: 800px;"\u003e')
            except Exception as e:
                print(f"Fallback PDF image extraction failed: {e}")

            # Insert any remaining tables that weren't positioned
            for tr in table_regions:
                if id(tr) not in inserted_tables:
                    try:
                        tab = tr['table']
                        df = tab.to_pandas()
                        html_parts.append('\u003ch4\u003eTable:\u003c/h4\u003e')
                        html_parts.append(df.to_html(index=False, classes="content-table").replace('class="dataframe content-table"', 'class="content-table"'))
                    except Exception as e:
                        print(f"Error rendering remaining table: {e}")

            html_parts.append('\u003c/div\u003e')

        html_parts.append('</div>')
        
        # [SCAN DETECTION] If total text blocks is very low relative to total pages, it's a scan.
        avg_text_per_page = total_text_blocks / len(doc) if len(doc) > 0 else 0
        if avg_text_per_page < 0.5: # Heuristic: less than 1 text block per 2 pages
             scan_warning = (
                 '<div class="note" style="background-color: #fee2e2; border: 2px solid #ef4444; color: #991b1b; padding: 20px; border-radius: 8px; margin-bottom: 25px;">'
                 '<strong>⚠️ ACCESSIBILITY WARNING: POTENTIAL SCANNED IMAGE</strong><br>'
                 'This PDF appears to be a scanned image of a document rather than a text-based file. '
                 'Screen readers will NOT be able to read this content. We have extracted images of the pages, '
                 'but we strongly recommend finding a text-based version or using an OCR tool (like Adobe Acrobat Pro or Microsoft Lens) '
                 'before converting to HTML.'
                 '</div>'
             )
             html_parts.insert(1, scan_warning)
             print(f"    [WARNING] PDF '{filename}' appears to be a scanned image (text blocks: {total_text_blocks}, pages: {len(doc)}).")

        full_content = "\n".join(html_parts)
        
        # Logging
        h_count = full_content.count("<h3") + full_content.count("<h2")
        img_count = full_content.count("<img")
        print(f"    [LOG] PDF Conversion: Detected {h_count} headers and {img_count} images.")
        
        s_filename = sanitize_filename(filename)
        output_path = os.path.join(output_dir, f"{s_filename}.html")
        
        _save_html(full_content, filename, pdf_path, output_path)
        return output_path, None

    except Exception as e:
        return None, f"PyMuPDF Error: {str(e)}"

def _convert_pdf_fallback(pdf_path):
    """Legacy text-only converter using pdfminer."""
    try:
        text = extract_text(pdf_path)
        paragraphs = text.split('\n\n')
        html_parts = []
        html_parts.append('<div class="pdf-content">')
        html_parts.append('<p class="note">⚠️ Text-Only Extraction (Images Missing).</p>')
        
        for p in paragraphs:
             clean_p = p.strip()
             if clean_p:
                  html_parts.append(f"<p>{clean_p}</p>")
        
        html_parts.append('</div>')
        full_content = "\n".join(html_parts)
        filename = os.path.splitext(os.path.basename(pdf_path))[0]
        s_filename = sanitize_filename(filename)
        output_path = os.path.join(os.path.dirname(pdf_path), f"{s_filename}.html")
        _save_html(full_content, filename, pdf_path, output_path)
        return output_path, None

    except Exception as e:
        return None, str(e)


def update_links_in_directory(directory, old_filename, new_filename):
    """
    Scans all HTML files in directory and replaces links using BeautifulSoup.
    e.g. <a href="syllabus.docx">Click Here</a> -> <a href="syllabus.html">Syllabus</a>
    Replaces underscores with spaces in link text.
    """
    count = 0
    old_base = os.path.basename(old_filename)
    new_base = os.path.basename(new_filename)
    
    # URL encoded version for comparison
    old_base_enc = old_base.replace(' ', '%20')
    
    # Generate new link text: replace underscores with spaces
    # e.g. Object_Oriented.html -> Object Oriented
    link_text_base = os.path.splitext(new_base)[0]
    new_link_text = link_text_base.replace('_', ' ').strip()
    
    # Handle live URLs vs local files
    if new_filename.startswith('http'):
        new_href = new_filename
    else:
        new_href = new_base.replace(' ', '%20')

    for root, dirs, files in os.walk(directory):
        for file in files:
            if file.endswith('.html'):
                filepath = os.path.join(root, file)
                try:
                    with open(filepath, 'r', encoding='utf-8') as f:
                        soup = BeautifulSoup(f.read(), 'html.parser')
                    
                    modified = False
                    # 1. Update Links (<a> tags)
                    for a in soup.find_all('a', href=True):
                        href = a['href']
                        # Standardize href for comparison
                        clean_href = urllib.parse.unquote(href).replace('\\', '/')
                        
                        if clean_href.endswith(old_base.replace('\\', '/')) or href == old_base_enc:
                            a['href'] = new_href
                            # Update link text automatically per user request
                            a.string = new_link_text
                            modified = True
                    
                    # 2. Update Images (<img> tags)
                    for img in soup.find_all('img', src=True):
                        src = img['src']
                        if src == old_base or src == old_base_enc:
                            img['src'] = new_href
                            modified = True

                    if modified:
                        with open(filepath, 'w', encoding='utf-8') as f:
                            f.write(str(soup))
                        count += 1
                except Exception as e:
                    print(f"Error updating links in {file}: {e}")
    return count


def unzip_course_package(zip_path, extract_to, log_func=None):
    """
    Extracts a Canvas Export (.imscc) or Zip file to the target directory.
    Renames .imscc to .zip internally if needed.
    """
    try:
        if not os.path.exists(extract_to):
            os.makedirs(extract_to)
            
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            members = zip_ref.namelist()
            total = len(members)
            for i, member in enumerate(members):
                # Check for stop request via log_func
                if log_func and hasattr(log_func, '__self__') and hasattr(log_func.__self__, 'stop_requested'):
                    if log_func.__self__.stop_requested:
                        return False, "Extraction stopped by user."
                
                zip_ref.extract(member, extract_to)
                
                if log_func and (i + 1) % 50 == 0:
                    log_func(f"   ... Extracted {i + 1}/{total} files...")

        return True, f"Success! Extracted to: {extract_to}"
    except Exception as e:
        return False, str(e)

def create_course_package(source_dir, output_path, log_func=None):
    """
    Zips the directory back into a .imscc file.
    Automatically excludes:
    - The originals archive folder
    - The output file itself (handles Windows case-insensitivity)
    - System/Dev folders like .git, venv, __pycache__
    """
    try:
        # Get absolute path of output to prevent zipping it into itself
        abs_output = os.path.normpath(os.path.abspath(output_path)).lower()
        
        # Folders to skip
        SKIP_DIRS = [ARCHIVE_FOLDER_NAME, '.git', 'venv', '.venv', '__pycache__', '.pytest_cache']
        
        file_count = 0
        total_files_added = 0
        
        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED, allowZip64=True) as zipf:
            for root, dirs, files in os.walk(source_dir):
                # Filter out skip directories
                # Modifying dirs in-place affects os.walk behavior
                dirs[:] = [d for d in dirs if d not in SKIP_DIRS]
                
                for file in files:
                    # Check for stop request via log_func
                    if log_func and hasattr(log_func, '__self__') and hasattr(log_func.__self__, 'stop_requested'):
                        if log_func.__self__.stop_requested:
                            # Close zip file and remove partial file if possible
                            # zipf is closed by 'with' block
                            return False, "Packaging stopped by user."

                    file_path = os.path.join(root, file)
                    abs_file = os.path.normpath(os.path.abspath(file_path)).lower()
                    
                    # [CRITICAL FIX] Skip the output .imscc file (Case-Insensitive for Windows)
                    if abs_file == abs_output:
                        continue

                    # Archive name should be relative to source_dir
                    arcname = os.path.relpath(file_path, source_dir)
                    zipf.write(file_path, arcname)
                    
                    file_count += 1
                    total_files_added += 1
                    
                    # Log progress every 50 files
                    if log_func and file_count >= 50:
                        log_func(f"   ... Added {total_files_added} files...")
                        file_count = 0
                        
        size_mb = os.path.getsize(output_path) / (1024 * 1024)
        return True, f"Created: {output_path} ({total_files_added} files, {size_mb:.2f} MB)"
    except Exception as e:
        return False, str(e)

def archive_source_file(file_path):
    """
    Moves an original source file to the archive folder.
    Returns the new path.
    """
    try:
        if not os.path.exists(file_path):
            return None
            
        dir_name = os.path.dirname(file_path)
        archive_dir = os.path.join(dir_name, ARCHIVE_FOLDER_NAME)
        
        if not os.path.exists(archive_dir):
            os.makedirs(archive_dir)
            
        new_path = os.path.join(archive_dir, os.path.basename(file_path))
        
        # If file already exists in archive, add a timestamp or just overwrite
        if os.path.exists(new_path):
            # For simplicity, we just move it. Shutil.move handles destination objects.
            pass
            
        shutil.move(file_path, new_path)
        return new_path
    except Exception as e:
        print(f"Error archiving {file_path}: {e}")
        return None

def update_manifest_resource(root_dir, old_rel_path, new_rel_path):
    """
    Updates imsmanifest.xml in the root_dir to reflect file changes.
    Replaces all occurrences of old_rel_path with new_rel_path in href attributes.
    """
    manifest_path = os.path.join(root_dir, 'imsmanifest.xml')
    if not os.path.exists(manifest_path):
        return False, "imsmanifest.xml not found"

    try:
        # Standardize paths to forward slashes for XML
        old_p = old_rel_path.replace("\\", "/").lower()
        new_p = new_rel_path.replace("\\", "/")

        with open(manifest_path, 'r', encoding='utf-8') as f:
            content = f.read()

        # Find all href="..." and replace if they match old_p (case-insensitive comparison)
        replacements = 0
        def repl_func(match):
            nonlocal replacements
            href_val = match.group(1)
            if href_val.replace("\\", "/").lower() == old_p:
                replacements += 1
                return f'href="{new_p}"'
            return match.group(0)

        new_content = re.sub(r'href="([^"]+)"', repl_func, content)

        if replacements > 0:
            with open(manifest_path, 'w', encoding='utf-8') as f:
                f.write(new_content)
            return True, f"Manifest Updated: {replacements} resource(s) synchronized."
        
        return False, "No matching entries found in imsmanifest.xml."
    except Exception as e:
        return False, f"Manifest update error: {str(e)}"

def run_janitor_cleanup(source_dir, log_func=None):
    """
    Scans the course for original source files (Word, PPT, PDF) that have been converted.
    Moves them to the archive folder to keep the Canvas course clean.
    """
    if log_func: log_func("🧹 Janitor: Tidying up original files...")
    
    extensions_to_clean = ['.docx', '.doc', '.pptx', '.ppt', '.xlsx', '.xls', '.pdf']
    cleaned_count = 0
    
    for root, dirs, files in os.walk(source_dir):
        # Don't clean files already in the archive or hidden folders
        if ARCHIVE_FOLDER_NAME in root or '.git' in root:
            continue
            
        for file in files:
            ext = os.path.splitext(file)[1].lower()
            if ext in extensions_to_clean:
                # [SAFETY FIX] Only archive if a converted version actually exists!
                # This prevents deleting files that the user chose NOT to convert.
                file_path = os.path.join(root, file)
                base_name = os.path.splitext(file)[0]
                html_version = os.path.join(root, base_name + ".html")
                
                if os.path.exists(html_version):
                    new_path = archive_source_file(file_path)
                    if new_path:
                        cleaned_count += 1
                else:
                    # Keep it in the course if not converted
                    pass
                    
    if log_func: log_func(f"🧹 Janitor: archived {cleaned_count} source files. (Safe for upload!)")
    return cleaned_count

